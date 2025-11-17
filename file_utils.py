import io
import json

def extract_text_from_docx(docx_bytes):
    from docx import Document
    f = io.BytesIO(docx_bytes)
    doc = Document(f)
    paragraphs = [p.text for p in doc.paragraphs]
    return "\n".join(paragraphs)

def extract_text_from_pptx(pptx_bytes):
    from pptx import Presentation
    f = io.BytesIO(pptx_bytes)
    pres = Presentation(f)
    slides = []
    for slide in pres.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        slides.append("\n".join(texts))
    full = "\n\n---\n\n".join(slides)
    return full, len(slides)

def extract_text_from_xlsx_csv(bytes_data, ftype):
    import pandas as pd
    f = io.BytesIO(bytes_data)
    if ftype == "csv":
        df = pd.read_csv(f, nrows=200)  # limit preview
    else:
        df = pd.read_excel(f, nrows=200)
    # convert to CSV-like preview text
    return df.head(50).to_csv(index=False)

def extract_text_from_text_or_code(bytes_data, encoding="utf-8"):
    try:
        return bytes_data.decode(encoding)
    except:
        return str(bytes_data[:1000])

def extract_text_from_ipynb(nb_bytes):
    import nbformat
    nb = nbformat.reads(nb_bytes.decode("utf-8"), as_version=4)
    texts = []
    for cell in nb.cells:
        if cell.cell_type == "markdown":
            texts.append(cell.source)
        elif cell.cell_type == "code":
            texts.append("```python\n" + cell.source + "\n```")
    return "\n\n".join(texts)
